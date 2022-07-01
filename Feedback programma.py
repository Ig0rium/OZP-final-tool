 

!python -m pip install python-pptx

from pptx import Presentation
#add the location of the presentatoin you want to analyse in the brackets at pr1 = Presentation
pr1 = Presentation()

#This function goes through every slide in the presentation and over every shape within a slide. If the shape is a text frame, it will read the words within and put these into a list. A dictionary is crerated with the number of the slide and the amount of words on that slide.

def count_words_per_slide (prs):
    wordcount_per_slide = {}
    slide_counter = 0
    for slide in prs.slides:
        slide_counter+=1
        total_words = 0
        #shapes are objects that contain information e.g., text, graphs, pictures
        for shape in slide.shapes:
            all_words = []
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    #runs are objects with information about the text
                    for run in paragraph.runs:
                        split_words = run.text.split()
                        all_words += split_words
                total_words += len(all_words)
                wordcount_per_slide.update({slide_counter:total_words})

    return (wordcount_per_slide)


def count_visuals_per_slide (prs):
    visuals_per_slide = {}
    slide_counter = 0
    for slide in prs.slides:
        slide_counter+=1
        total_visuals = []
        for shape in slide.shapes:
            if "Afbeelding" in shape.name:
                total_visuals.append(shape.name)
                visuals_per_slide.update({slide_counter:len(total_visuals)})
            if "Image" in shape.name:
                total_visuals.append(shape.name)
                visuals_per_slide.update({slide_counter:len(total_visuals)})
            else:
                visuals_per_slide.update({slide_counter:len(total_visuals)})
    return (visuals_per_slide)

def words_feedback(dict):
    slides15words = []
    slides40words = []
    for key in dict:
        if 40 > dict[key] >= 15:
            slides15words.append(key)
        if dict[key] >= 40:
            slides40words.append(key)
    print("The following slides contain more than 15 and less than 40 words. Check to see if all text is necessary:",*slides15words, sep=', ' )
    print("The following slides contain more than 40 words. Strongly consider deleting some text in these slides:",*slides40words, sep=', ' )

def visuals_feedback(dict):
    slides_without_visuals = []
    for key in dict:
        if dict[key] < 1:
            slides_without_visuals.append(key)
    print("The following slides contain no visuals. Consider adding some visuals to these slides:", *slides_without_visuals, sep=', ' )
   
def good_slide_feedback(words,visuals):
    good_slides = []
    for key in words:
        if words[key] < 15:
            if visuals[key] > 0:
                good_slides.append(key)
    print("The following slides contain less than 15 words and at least 1 visual, which means they are likely more effective than other slides:",*good_slides, sep=', ')


visuals_feedback(count_visuals_per_slide(pr1))
words_feedback(count_words_per_slide(pr1))
good_slide_feedback(count_words_per_slide(pr1),count_visuals_per_slide(pr1))
